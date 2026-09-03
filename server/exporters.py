# -*- coding: utf-8 -*-
"""
PDF 에서 다른 형식으로 내보내는 변환기들.

한글·오피스 문서를 PDF 로 바꾸는 converter.py 와 반대 방향이다. 이쪽은 COM 을
쓰지 않지만(PDF/A 의 Ghostscript 제외) 무거운 작업이라 같은 큐를 태워 순서대로
처리한다.
"""
from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from typing import Optional

from converter import ConversionError


# ---------------------------------------------------------------------------
# PDF -> Word (.docx)
# ---------------------------------------------------------------------------

def pdf_to_docx(src: str, dst: str, progress=None) -> None:
    """
    레이아웃을 최대한 살려 Word 문서로 바꾼다.

    MS Word 를 쓰지 않는다. Word COM 으로 PDF 를 열면 "PDF 를 편집 가능한 문서로
    변환합니다" 확인 대화상자가 떠서 무인 실행이 멈추고, 이를 없애려면 사용자의
    Office 설정을 바꿔야 한다. 남의 설정을 건드리지 않으려고 순수 파이썬 경로를 쓴다.
    """
    try:
        from pdf2docx import Converter
    except ImportError as e:
        raise ConversionError(
            "PDF→Word 변환에 필요한 구성요소가 없습니다. start.bat 을 다시 실행해 주세요."
        ) from e

    cv = Converter(src)
    try:
        cv.convert(dst, start=0, end=None)
    except Exception as e:
        raise ConversionError(
            "이 PDF는 Word 문서로 바꾸지 못했습니다. "
            "글자가 없는 스캔 이미지 PDF라면 먼저 OCR을 돌린 뒤 다시 시도해 보세요."
        ) from e
    finally:
        cv.close()

    if not os.path.exists(dst):
        raise ConversionError("Word 문서가 만들어지지 않았습니다.")


# ---------------------------------------------------------------------------
# PDF -> Excel (.xlsx)
# ---------------------------------------------------------------------------

def pdf_to_xlsx(src: str, dst: str, progress=None) -> None:
    """
    PDF 안의 표를 찾아 시트로 옮긴다.

    표가 하나도 없으면 실패로 처리한다. 빈 통합문서를 주는 것보다 왜 안 됐는지
    알려주는 편이 낫다.
    """
    import pymupdf
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter

    doc = pymupdf.open(src)
    wb = Workbook()
    wb.remove(wb.active)
    found = 0

    try:
        for pno in range(doc.page_count):
            if progress:
                progress(pno, doc.page_count, f"{pno + 1}쪽에서 표 찾는 중")
            page = doc[pno]
            try:
                tables = page.find_tables()
            except Exception:
                continue
            for tno, table in enumerate(tables.tables, start=1):
                rows = table.extract()
                if not rows:
                    continue
                found += 1
                title = f"{pno + 1}쪽" + (f"-{tno}" if tno > 1 else "")
                ws = wb.create_sheet(title[:31])
                widths = {}
                for r in rows:
                    cleaned = [("" if c is None else str(c).replace("\n", " ").strip()) for c in r]
                    ws.append(cleaned)
                    for i, val in enumerate(cleaned, start=1):
                        widths[i] = min(60, max(widths.get(i, 10), len(val) + 2))
                for i, w in widths.items():
                    ws.column_dimensions[get_column_letter(i)].width = w
    finally:
        doc.close()

    if not found:
        raise ConversionError(
            "이 PDF에서 표를 찾지 못했습니다. 표 형태가 뚜렷한 문서에서만 동작합니다. "
            "선이 없는 표나 스캔한 이미지는 인식하지 못합니다."
        )

    wb.save(dst)


# ---------------------------------------------------------------------------
# PDF -> PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def pdf_to_pptx(src: str, dst: str, dpi: int = 150, progress=None) -> None:
    """
    각 쪽을 슬라이드 한 장의 배경 그림으로 넣는다.

    PDF 에는 슬라이드라는 개념이 없어서 도형이나 텍스트 상자를 되살릴 방법이 없다.
    쪽 모양을 그대로 보존하는 쪽을 택했다. 화면에 띄우거나 인쇄하는 용도로는 충분하다.
    """
    import pymupdf
    from pptx import Presentation
    from pptx.util import Emu

    doc = pymupdf.open(src)
    if doc.page_count == 0:
        doc.close()
        raise ConversionError("페이지가 없는 PDF입니다.")

    prs = Presentation()
    blank = prs.slide_layouts[6]           # 빈 레이아웃
    tmpdir = tempfile.mkdtemp(prefix="pptx_")

    try:
        # 첫 쪽 비율에 슬라이드 크기를 맞춘다(가로 문서·세로 문서 모두 자연스럽게).
        first = doc[0].rect
        ratio = first.height / first.width if first.width else 0.75
        prs.slide_width = Emu(9144000)                       # 10인치
        prs.slide_height = Emu(int(9144000 * ratio))

        for pno in range(doc.page_count):
            if progress:
                progress(pno, doc.page_count, f"{pno + 1}쪽 슬라이드로 변환 중")
            pix = doc[pno].get_pixmap(dpi=dpi)
            img_path = os.path.join(tmpdir, f"p{pno:04d}.png")
            pix.save(img_path)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(img_path, 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        prs.save(dst)
    finally:
        doc.close()
        shutil.rmtree(tmpdir, ignore_errors=True)


# ---------------------------------------------------------------------------
# HTML / 웹페이지 -> PDF
# ---------------------------------------------------------------------------

_BROWSER_CANDIDATES = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]


def find_browser() -> Optional[str]:
    """설치된 Chrome 또는 Edge 를 찾는다. 둘 다 PDF 인쇄를 지원한다."""
    for path in _BROWSER_CANDIDATES:
        if os.path.exists(path):
            return path
    for name in ("chrome", "msedge"):
        found = shutil.which(name)
        if found:
            return found
    return None


def html_to_pdf(url: str, dst: str, *, wait_ms: int = 4000) -> None:
    """
    웹페이지나 HTML 파일을 PDF 로 인쇄한다.

    Playwright 를 따로 받지 않고 이미 깔려 있는 Chrome/Edge 의 헤드리스 모드를 쓴다.
    150MB 짜리 브라우저를 또 내려받을 이유가 없다.

    용지 크기와 방향은 여기서 지정할 수 없다. 헤드리스 Chrome 의 --print-to-pdf 는
    --landscape 같은 인자를 받지 않고 문서의 @page CSS 만 따른다(실측 확인).
    업로드한 HTML 파일이라면 html_file_to_pdf 로 @page 를 넣어 조절할 수 있다.
    """
    browser = find_browser()
    if not browser:
        raise ConversionError(
            "Chrome 또는 Microsoft Edge를 찾을 수 없습니다. 둘 중 하나가 필요합니다."
        )

    # 사용자가 쓰고 있는 브라우저 창에 끼어들지 않도록 임시 프로필로 따로 띄운다.
    profile = tempfile.mkdtemp(prefix="hanmoa_browser_")
    try:
        cmd = [
            browser,
            "--headless=new",
            "--disable-gpu",
            "--no-first-run",
            "--no-default-browser-check",
            "--user-data-dir=" + profile,
            "--virtual-time-budget=" + str(wait_ms),
            "--no-pdf-header-footer",
            "--print-to-pdf=" + dst,
            url,
        ]
        proc = subprocess.run(cmd, capture_output=True, timeout=180)
        if not os.path.exists(dst) or os.path.getsize(dst) == 0:
            detail = proc.stderr.decode("utf-8", errors="ignore")[-300:].strip()
            raise ConversionError(
                "페이지를 PDF로 만들지 못했습니다. 주소가 맞는지, 로그인이 필요한 "
                "페이지는 아닌지 확인해 주세요." + (("\n(" + detail + ")") if detail else "")
            )
    except subprocess.TimeoutExpired as e:
        raise ConversionError("페이지를 불러오는 데 너무 오래 걸려 중단했습니다.") from e
    finally:
        shutil.rmtree(profile, ignore_errors=True)


PAPER_SIZES = {
    "A4": "A4", "A3": "A3", "Letter": "Letter", "Legal": "Legal",
}


def html_file_to_pdf(src: str, dst: str, *, paper: str = "A4",
                     landscape: bool = False, margin_mm: int = 12,
                     wait_ms: int = 3000) -> None:
    """
    업로드한 HTML 파일을 PDF 로 만든다. 용지 크기와 방향을 지정할 수 있다.

    브라우저는 문서의 @page 규칙만 보고 용지를 정하므로, 원본 옆에 사본을 만들어
    @page 규칙을 앞에 끼워 넣는다. 사본을 같은 폴더에 두어야 이미지·CSS 같은
    상대 경로 자원이 그대로 따라온다.
    """
    size = PAPER_SIZES.get(paper, "A4")
    orientation = " landscape" if landscape else ""
    rule = (
        "<style>@page{size:" + size + orientation + ";"
        "margin:" + str(int(margin_mm)) + "mm;}</style>"
    )

    raw = open(src, "rb").read()
    for codec in ("utf-8", "cp949", "latin-1"):
        try:
            text = raw.decode(codec)
            break
        except UnicodeDecodeError:
            continue

    # 문자 인코딩 선언보다 뒤에 넣어야 한글이 깨지지 않는다.
    lowered = text.lower()
    pos = lowered.find("<head>")
    if pos >= 0:
        text = text[:pos + 6] + '<meta charset="utf-8">' + rule + text[pos + 6:]
    else:
        text = '<meta charset="utf-8">' + rule + text

    staged = os.path.join(os.path.dirname(src), "_print_" + os.path.basename(src))
    if not staged.lower().endswith((".html", ".htm")):
        staged += ".html"
    with open(staged, "w", encoding="utf-8") as fh:
        fh.write(text)

    try:
        html_to_pdf("file:///" + staged.replace("\\", "/"), dst, wait_ms=wait_ms)
    finally:
        try:
            os.remove(staged)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# PDF -> PDF/A (장기 보존용)
# ---------------------------------------------------------------------------

_GS_CANDIDATES = ["gswin64c", "gswin32c", "gs"]


def find_ghostscript() -> Optional[str]:
    """Ghostscript 실행 파일을 찾는다. PATH 와 기본 설치 위치를 본다."""
    for name in _GS_CANDIDATES:
        found = shutil.which(name)
        if found:
            return found
    for root in (r"C:\Program Files\gs", r"C:\Program Files (x86)\gs"):
        if not os.path.isdir(root):
            continue
        for entry in sorted(os.listdir(root), reverse=True):     # 최신 버전 우선
            for exe in ("gswin64c.exe", "gswin32c.exe"):
                path = os.path.join(root, entry, "bin", exe)
                if os.path.exists(path):
                    return path
    return None


def pdf_to_pdfa(src: str, dst: str, *, level: str = "2") -> None:
    """
    PDF/A 로 변환한다. 글꼴을 문서 안에 모두 넣어 오랜 시간이 지나도 같은 모양으로
    열리게 만드는 보존용 형식이다.

    Ghostscript 가 필요하다. 폰트 임베딩과 색 공간 변환을 제대로 하는 도구가
    사실상 이것뿐이라 대체 구현을 두지 않았다.
    """
    gs = find_ghostscript()
    if not gs:
        raise ConversionError(
            "PDF/A 변환에는 Ghostscript가 필요한데 설치되어 있지 않습니다. "
            "https://ghostscript.com/releases/gsdnld.html 에서 내려받아 설치한 뒤 "
            "start.bat 을 다시 실행해 주세요."
        )

    cmd = [
        gs, "-dBATCH", "-dNOPAUSE", "-dQUIET", "-dSAFER",
        "-dPDFA=" + str(level), "-dPDFACompatibilityPolicy=1",
        "-sDEVICE=pdfwrite",
        "-sColorConversionStrategy=UseDeviceIndependentColor",
        "-dEmbedAllFonts=true", "-dSubsetFonts=true",
        "-sOutputFile=" + dst, src,
    ]
    proc = subprocess.run(cmd, capture_output=True, timeout=300)
    if not os.path.exists(dst) or os.path.getsize(dst) == 0:
        detail = proc.stderr.decode("utf-8", errors="ignore")[-300:].strip()
        raise ConversionError("PDF/A 변환에 실패했습니다." + (("\n(" + detail + ")") if detail else ""))


# ---------------------------------------------------------------------------
# 손상된 PDF 복구
# ---------------------------------------------------------------------------

def repair_pdf(src: str, dst: str) -> None:
    """
    깨진 PDF 에서 읽을 수 있는 부분을 건져 다시 쓴다.

    MuPDF 는 상호참조표가 망가진 파일도 객체를 훑어 다시 세우기 때문에,
    다른 뷰어가 못 여는 파일이 열리는 경우가 있다.
    """
    import pymupdf

    try:
        doc = pymupdf.open(src)
    except Exception as e:
        raise ConversionError(
            "파일을 전혀 읽을 수 없습니다. PDF가 아니거나 손상 정도가 너무 큽니다."
        ) from e

    try:
        if doc.page_count == 0:
            raise ConversionError("살릴 수 있는 페이지를 찾지 못했습니다.")
        # garbage=4: 안 쓰는 객체 정리, clean=True: 구조 재작성
        doc.save(dst, garbage=4, clean=True, deflate=True)
    finally:
        doc.close()
